"""
PPT生成服务 - 4套模板 + 5种循环版式 + 每子要点一页
"""
import os
import uuid
from io import BytesIO
from typing import List, Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import requests
from loguru import logger


COLOR_THEMES = {
    "academic": {
        "name": "学术蓝", "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title": RGBColor(0x1A, 0x3C, 0x6D), "text": RGBColor(0x33, 0x33, 0x33),
        "accent": RGBColor(0x2B, 0x7A, 0xC2), "line": RGBColor(0x1A, 0x3C, 0x6D),
        "light_bg": RGBColor(0xF0, 0xF4, 0xFA),
    },
    "fresh": {
        "name": "清新绿", "bg": RGBColor(0xFA, 0xFC, 0xFF),
        "title": RGBColor(0x2D, 0x8C, 0x5A), "text": RGBColor(0x44, 0x44, 0x44),
        "accent": RGBColor(0x5B, 0xB3, 0x81), "line": RGBColor(0x2D, 0x8C, 0x5A),
        "light_bg": RGBColor(0xF0, 0xFA, 0xF4),
    },
    "tech": {
        "name": "科技紫", "bg": RGBColor(0x0D, 0x11, 0x17),
        "title": RGBColor(0x71, 0x70, 0xFF), "text": RGBColor(0xE0, 0xE0, 0xE0),
        "accent": RGBColor(0x9D, 0x71, 0xFF), "line": RGBColor(0x71, 0x70, 0xFF),
        "light_bg": RGBColor(0x18, 0x1C, 0x24),
    },
    "warm": {
        "name": "暖橙", "bg": RGBColor(0xFF, 0xFD, 0xF8),
        "title": RGBColor(0xC0, 0x5A, 0x20), "text": RGBColor(0x4A, 0x3A, 0x30),
        "accent": RGBColor(0xE8, 0x7A, 0x3C), "line": RGBColor(0xC0, 0x5A, 0x20),
        "light_bg": RGBColor(0xFF, 0xF5, 0xEC),
    },
}


class PPTGenerator:
    def __init__(self, style_name: str = "academic"):
        self.c = COLOR_THEMES.get(style_name, COLOR_THEMES["academic"])
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.page_num = 1
        self._font = "Microsoft YaHei"

    def _bg(self, slide, color=None):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color or self.c["bg"]

    def _line(self, slide, l, t, w):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(3))
        s.fill.solid()
        s.fill.fore_color.rgb = self.c["line"]
        s.line.fill.background()

    def _textbox(self, slide, l, t, w, h, text, size=16, bold=False, color=None, align=PP_ALIGN.LEFT, spacing=8):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color or self.c["text"]
        p.font.name = self._font
        p.alignment = align
        p.space_after = Pt(spacing)
        return tf

    def _multiline(self, tf, lines, size=15, color=None, bold_first=False):
        for i, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(size)
            p.font.color.rgb = color or self.c["text"]
            p.font.name = self._font
            p.space_after = Pt(8)
            if bold_first and i == 0:
                p.font.bold = True

    def _page_number(self, slide):
        tb = slide.shapes.add_textbox(Inches(12.2), Inches(7.0), Inches(0.9), Inches(0.35))
        p = tb.text_frame.paragraphs[0]
        p.text = str(self.page_num)
        p.font.size = Pt(10)
        p.font.color.rgb = self.c["accent"]
        p.font.name = self._font
        p.alignment = PP_ALIGN.RIGHT
        self.page_num += 1

    def _download_image(self, url):
        try:
            resp = requests.get(url, timeout=30)
            if resp.status_code == 200:
                return BytesIO(resp.content)
        except Exception:
            pass
        return None

    # ====== 六种版式 ======

    def slide_cover(self, title, subtitle=""):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide)
        self._textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
                      title, size=40, bold=True, color=self.c["title"], align=PP_ALIGN.CENTER)
        self._line(slide, Inches(4.5), Inches(3.5), Inches(4.3))
        if subtitle:
            self._textbox(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.8),
                          subtitle, size=16, color=self.c["accent"], align=PP_ALIGN.CENTER)
        self._page_number(slide)

    def slide_toc(self, sections):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide)
        self._textbox(slide, Inches(1.2), Inches(0.5), Inches(10), Inches(0.7),
                      "课程目录", size=30, bold=True, color=self.c["title"])
        self._line(slide, Inches(1.2), Inches(1.2), Inches(4))
        y = Inches(1.5)
        for i, sec in enumerate(sections):
            dur = sec.get("duration", 0)
            self._textbox(slide, Inches(1.5), y, Inches(10), Inches(0.4),
                          f"{i+1:02d}  {sec.get('title', '')}  ({dur}分钟)",
                          size=15, color=self.c["text"])
            y += Inches(0.42)
        self._page_number(slide)

    def _add_section_title_bar(self, slide, section_title, section_index):
        """顶部章节标题条（全宽深色条）"""
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.c["title"]
        bar.line.fill.background()
        prefix = f"0{section_index + 1}" if section_index < 9 else str(section_index + 1)
        self._textbox(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.7),
                      f"{prefix}  {section_title}", size=24, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF))

    def slide_full_text(self, section_title, point_title, content, section_index):
        """版式1: 全宽文字页（左标题 + 右要点）"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide)
        self._add_section_title_bar(slide, section_title, section_index)
        self._textbox(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.6),
                      point_title, size=22, bold=True, color=self.c["title"])
        self._line(slide, Inches(0.8), Inches(1.95), Inches(3))
        lines = [l.strip() for l in content.split('\n') if l.strip()]
        tf = self._textbox(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.5),
                           "", size=15, color=self.c["text"])
        self._multiline(tf, lines, size=15)
        self._page_number(slide)

    def slide_image_left(self, section_title, point_title, content, image_url, section_index):
        """版式2: 左图右文"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide)
        self._add_section_title_bar(slide, section_title, section_index)
        self._textbox(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.6),
                      point_title, size=20, bold=True, color=self.c["title"])
        self._line(slide, Inches(0.8), Inches(1.9), Inches(2.5))
        lines = [l.strip() for l in content.split('\n') if l.strip()]
        tf = self._textbox(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(4.2),
                           "", size=14, color=self.c["text"])
        self._multiline(tf, lines, size=14)
        img_data = self._download_image(image_url)
        if img_data:
            slide.shapes.add_picture(img_data, Inches(6.8), Inches(1.5), width=Inches(5.8))
        self._page_number(slide)

    def slide_image_right(self, section_title, point_title, content, image_url, section_index):
        """版式3: 左文右图"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide)
        self._add_section_title_bar(slide, section_title, section_index)
        self._textbox(slide, Inches(0.8), Inches(1.4), Inches(7), Inches(0.6),
                      point_title, size=20, bold=True, color=self.c["title"])
        self._line(slide, Inches(0.8), Inches(1.9), Inches(2.5))
        lines = [l.strip() for l in content.split('\n') if l.strip()]
        tf = self._textbox(slide, Inches(0.8), Inches(2.1), Inches(7), Inches(4.2),
                           "", size=14, color=self.c["text"])
        self._multiline(tf, lines, size=14)
        img_data = self._download_image(image_url)
        if img_data:
            slide.shapes.add_picture(img_data, Inches(8.3), Inches(1.5), width=Inches(4.3))
        self._page_number(slide)

    def slide_bullet_list(self, section_title, point_title, content, section_index):
        """版式4: 卡片式要点列表"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide)
        self._add_section_title_bar(slide, section_title, section_index)
        self._textbox(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
                      point_title, size=22, bold=True, color=self.c["title"])
        self._line(slide, Inches(0.8), Inches(1.8), Inches(3))
        lines = [l.strip() for l in content.split('\n') if l.strip()]
        y = Inches(2.1)
        for i, line in enumerate(lines[:8]):
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(1.0), y, Inches(11), Inches(0.55))
            card.fill.solid()
            card.fill.fore_color.rgb = self.c["light_bg"]
            card.line.fill.background()
            tf = card.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.3)
            tf.margin_top = Pt(6)
            p = tf.paragraphs[0]
            p.text = f"   {line}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.c["text"]
            p.font.name = self._font
            y += Inches(0.62)
        self._page_number(slide)

    def slide_summary(self, section_title, point_title, content, section_index):
        """版式5: 总结页（居中大标题 + 关键点）"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide)
        self._textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1),
                      point_title, size=30, bold=True, color=self.c["title"], align=PP_ALIGN.CENTER)
        self._line(slide, Inches(5), Inches(2.5), Inches(3.3))
        lines = [l.strip() for l in content.split('\n') if l.strip()]
        tf = self._textbox(slide, Inches(2.0), Inches(2.8), Inches(9), Inches(4),
                           "", size=16, color=self.c["text"], align=PP_ALIGN.CENTER)
        self._multiline(tf, lines, size=16, color=self.c["text"])
        self._page_number(slide)

    layout_funcs = ["slide_full_text", "slide_bullet_list", "slide_full_text", "slide_image_right", "slide_bullet_list"]
    # 有配图时优先图文版式
    image_layouts = ["slide_image_right", "slide_image_left", "slide_bullet_list", "slide_image_right", "slide_full_text"]

    def generate(self, lesson_data: Dict, output_path: str = None) -> str:
        title = lesson_data.get("title", "AI 生成教案")
        sections = lesson_data.get("sections", [])

        self.slide_cover(title, "AI 教学助手自动生成 | 共约 {} 页".format(
            sum(len(s.get("sub_points", [s.get("title","")])) for s in sections) + 3))
        self.slide_toc(sections)

        for si, section in enumerate(sections):
            sub_points = section.get("sub_points", [])
            if not sub_points:
                sub_points = [section.get("title", "本节内容")]
            section_content = section.get("content", "")
            section_image = section.get("image_url")
            content_lines = section_content.split('\n')

            layouts = self.image_layouts if section_image else self.layout_funcs

            for pi, sp in enumerate(sub_points):
                chunk_size = max(1, len(content_lines) // len(sub_points))
                start = pi * chunk_size
                end = start + chunk_size if pi < len(sub_points) - 1 else len(content_lines)
                point_content = '\n'.join(content_lines[start:end]) if content_lines else str(sp)

                layout_idx = (si * len(sub_points) + pi) % len(layouts)
                layout = layouts[layout_idx]

                if layout in ("slide_image_right", "slide_image_left") and section_image:
                    getattr(self, layout)(section.get("title", ""), str(sp), point_content, section_image, si)
                else:
                    if layout in ("slide_image_right", "slide_image_left"):
                        layout = "slide_bullet_list"
                    getattr(self, layout)(section.get("title", ""), str(sp), point_content, si)

        self.slide_summary("课程总结", "感谢聆听", "本节课主要内容已全部呈现\n祝教学顺利！", -1)

        if output_path is None:
            lesson_id = lesson_data.get("lesson_id", uuid.uuid4().hex[:8])
            output_dir = os.path.join(os.path.dirname(__file__), "..", "data", "lessons")
            os.makedirs(output_dir, exist_ok=True)
            output_path = os.path.join(output_dir, f"{lesson_id}.pptx")

        self.prs.save(output_path)
        logger.info(f"PPT 已生成: {output_path} ({self.page_num - 1} 页)")
        return output_path


def generate_ppt_from_lesson(lesson_data: Dict, style: str = "academic") -> str:
    generator = PPTGenerator(style_name=style)
    return generator.generate(lesson_data)
